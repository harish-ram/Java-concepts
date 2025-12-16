import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor

BASE = Path(__file__).resolve().parents[1]
MD_FILE = BASE / 'DOCS' / 'EXAMPLES_SYNTAX_OOP.md'
OUT_FILE = BASE / 'DOCS' / 'EXAMPLES_SYNTAX_OOP.pptx'


def add_title_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_box = slide.shapes.title
    title_box.text = title
    if subtitle:
        tx = slide.placeholders[1]
        tx.text = subtitle


def add_content_slide(prs, title, paragraphs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for p in paragraphs:
        p = p.rstrip('\n')
        if p.startswith('```') and p.endswith('```'):
            # code block in single line (rare)
            code = p[3:-3]
            para = body.add_paragraph()
            run = para.add_run()
            run.text = code
            run.font.name = 'Consolas'
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        elif p.startswith('```'):
            # start of code block
            code_lines = []
            # remove leading ```
            continue
        else:
            para = body.add_paragraph()
            para.text = p
            para.level = 0


# Better parser that preserves code blocks

def parse_markdown(md_text):
    sections = []
    cur_title = None
    cur_lines = []
    in_code = False
    code_buf = []

    for line in md_text.splitlines():
        if line.startswith('## '):
            if cur_title:
                if code_buf:
                    cur_lines.append('```' + '\n'.join(code_buf) + '```')
                    code_buf = []
                sections.append((cur_title, cur_lines))
            cur_title = line[3:].strip()
            cur_lines = []
            in_code = False
        elif line.startswith('```'):
            if in_code:
                # close code
                in_code = False
                cur_lines.append('```' + '\n'.join(code_buf) + '```')
                code_buf = []
            else:
                in_code = True
                code_buf = []
        else:
            if in_code:
                code_buf.append(line)
            else:
                if line.strip() == '':
                    continue
                # remove emoji bullets or file refs that duplicate
                cur_lines.append(line)
    if cur_title:
        if code_buf:
            cur_lines.append('```' + '\n'.join(code_buf) + '```')
        sections.append((cur_title, cur_lines))
    return sections


def add_slide_from_section(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for line in lines:
        if line.startswith('```') and line.endswith('```'):
            code = line[3:-3]
            para = tf.add_paragraph()
            run = para.add_run()
            run.text = code
            run.font.name = 'Consolas'
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        else:
            # normal text: break into shorter paragraphs if long
            for chunk in wrap_text(line, 300):
                p = tf.add_paragraph()
                p.text = chunk
                p.level = 0


def wrap_text(text, n):
    # naive wrap: split by sentence-ish boundaries
    lines = []
    while len(text) > n:
        idx = text.rfind(' ', 0, n)
        if idx == -1:
            idx = n
        lines.append(text[:idx])
        text = text[idx+1:]
    if text:
        lines.append(text)
    return lines


def main():
    md = MD_FILE.read_text(encoding='utf-8')
    prs = Presentation()
    # Title slide
    add_title_slide(prs, 'Java Examples: Syntax, Data Types & OOP', 'From DOCS/EXAMPLES_SYNTAX_OOP.md')
    sections = parse_markdown(md)
    for title, lines in sections:
        add_slide_from_section(prs, title, lines)
    prs.save(OUT_FILE)
    print(f'Wrote: {OUT_FILE}')

if __name__ == '__main__':
    main()
