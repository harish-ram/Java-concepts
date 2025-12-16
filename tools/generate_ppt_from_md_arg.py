import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

USAGE = 'Usage: python generate_ppt_from_md_arg.py <path/to/file.md>'


def add_title_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle


def parse_markdown(md_text):
    sections = []
    cur_title = None
    cur_lines = []
    in_code = False
    code_buf = []

    for line in md_text.splitlines():
        if line.startswith('## '):
            if cur_title:
                if code_buf:
                    cur_lines.append('```' + '\n'.join(code_buf) + '```')
                    code_buf = []
                sections.append((cur_title, cur_lines))
            cur_title = line[3:].strip()
            cur_lines = []
            in_code = False
        elif line.startswith('```'):
            if in_code:
                in_code = False
                cur_lines.append('```' + '\n'.join(code_buf) + '```')
                code_buf = []
            else:
                in_code = True
                code_buf = []
        else:
            if in_code:
                code_buf.append(line)
            else:
                if line.strip() == '':
                    continue
                cur_lines.append(line)
    if cur_title:
        if code_buf:
            cur_lines.append('```' + '\n'.join(code_buf) + '```')
        sections.append((cur_title, cur_lines))
    return sections


def add_slide_from_section(prs, title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for line in lines:
        if line.startswith('```') and line.endswith('```'):
            code = line[3:-3]
            para = tf.add_paragraph()
            run = para.add_run()
            run.text = code
            run.font.name = 'Consolas'
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        else:
            # normal text chunk
            p = tf.add_paragraph()
            p.text = line
            p.level = 0


def main():
    if len(sys.argv) != 2:
        print(USAGE)
        sys.exit(1)
    md_path = Path(sys.argv[1])
    if not md_path.exists():
        print(f'File not found: {md_path}')
        sys.exit(1)
    md_text = md_path.read_text(encoding='utf-8')
    prs = Presentation()
    add_title_slide(prs, md_path.stem.replace('_',' '), f'From {md_path.name}')
    sections = parse_markdown(md_text)
    for title, lines in sections:
        add_slide_from_section(prs, title, lines)
    out_path = md_path.with_suffix('.pptx')
    prs.save(out_path)
    print(f'Wrote: {out_path}')

if __name__ == '__main__':
    main()
